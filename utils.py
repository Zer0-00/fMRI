import torch
import numpy as np
from scipy import signal
from scipy import io
from torchaudio.functional import filtfilt
import os
import copy
import pptx
import configs


def read_bfiles(fileDir, type_data='b'):
    """read byte files of int8 and return a list of data

    Args:
        filename (str): directory of file to read
        type_data (str): type of the byte data('b' for int8 and 'h' for int16)
    Returns:
        list: data
    """
    from struct import unpack    
    data = []
    if type_data == 'b':
        len_bdata = 1
    elif type_data == 'h':
        len_bdata = 2
    with open(fileDir, 'rb') as f:
        while True:
            bdata = f.read(len_bdata)
            if len(bdata) < len_bdata:
                break
            bdata = unpack(type_data, bdata)
            data.append(bdata[0])
            
    return data

def fmask(fdata, mask):
    print(fdata.shape, mask.shape)
    device = fdata.device
    data_squeeze = torch.squeeze(fdata[:,:,:,1])
    assert data_squeeze.shape == mask.shape, "mask shape is not equal to data size"
    mask = mask.cpu().numpy()
    mask = torch.Tensor(mask.reshape(-1, order='F')).to(device)
    idx = torch.where(mask)[0].to(device)
    fdata = fdata.cpu().numpy()
    fdata_reshape = torch.Tensor(np.reshape(fdata, (-1, fdata.shape[3]), order='F')).to(device)
    masked_fdata = fdata_reshape[idx, :]
    return masked_fdata

def regress(y, X):
    mata = torch.mm(X.transpose(0,1), X)
    matb = torch.mm(X.transpose(0,1), y)
    if mata.shape[0]>1 or mata.shape[1]>1:
        Beta = torch.linalg.solve(mata, matb)
    else:
        Beta = matb/mata
    Residual = y - torch.mm(X,Beta)
    
    return Beta, Residual

class filter():
    def __init__(self, shape, fs, fc, order, device):
        """_summary_

        Args:
            shape (str): the type of the filter, choose from[‘lowpass’, ‘highpass’, ‘bandpass’, ‘bandstop’]
            fs (float/int): sampling rate
            fc (float/int): cut-off frequency
            order (int): order of the filter
        """
        
        self.epsilon = 1e-6
        self.fc = copy.deepcopy(fc)
        if len(self.fc) == 1:
            if self.fc == fs /2 :
                self.fc = self.fc - self.epsilon
        elif len(self.fc) == 2:
            if self.fc[1] == fs /2 :
                self.fc[1] = self.fc[1] - self.epsilon

                
        b, a = signal.butter(order, self.fc, btype=shape, fs=fs)
        self.b = torch.Tensor(b).to(device)
        self.a = torch.Tensor(a).to(device)
        
    def filter(self, data):
        """filter the data

        Args:
            data (torch.Tensor): input data

        Returns:
            torch.Tensor: filtered data
        """
        #normalize
        
        scale = torch.max(torch.abs(data))
        #norm to [-1,1]
        data_norm = (data/scale)
        
        data_filtered = filtfilt(data_norm, self.a, self.b, clamp=False)
        
        #denormalize
        data_denormed = data_filtered * scale
        #compare with scipy
        
        return data_denormed
        
def create_folder(folder_dir):
    if not os.path.exists(folder_dir):
        os.makedirs(folder_dir)

def save_h5(file_name, data, dtype=None):
    import h5py
    with h5py.File(file_name, 'w') as f:
        dset = f.create_dataset('data', data=data, dtype=dtype)
        
def save_mat(file_name, data):
    io.savemat(file_name, {'data':data})

def add_text(slide, text, x, y, width, height):
    textbox = slide.shapes.add_textbox(x, y, width, height)
    tf = textbox.text_frame
    tf.text = text
        
    return textbox

def save_pptx(configs:configs.configuration):
    
    
    #setting ppt name 
    #Add the specific part of the directory name based on the filter
    common_dir = configs.ppt_name + '_Nofilter' if filter == 0.5 else configs.ppt_name

    #Add the specific part of the directory name based on the GSindex
    common_dir = common_dir + '_NGSR' if configs.GSindex == 0 else common_dir
    
    if configs.Taskindex == 1:
    #Add the specific part of the directory name based on the AtlasbasedROIregres
        common_dir = common_dir + '_TaskRegres' if configs.AtlasbasedROIregres == 1 else common_dir + '_TaskRegres36pixels'

    #Add the specific part of the directory name based on the Motionregres
    common_dir = common_dir + '_MotionReg' if configs.Motionregres == 1 else common_dir
    
    #Add the specific part of the directory name based on the bandwidth
    ppt_name = common_dir + '_{}-{}Hz'.format(configs.fcs[configs.bandwidth][0], configs.fcs[configs.bandwidth][1]) if configs.bandwidth in [0,1,2,3,4] else common_dir
    
    #Add the specific part of the directory name based on the smoothing
    common_dir = common_dir + '_smooth' if configs.gaussian_sigma > 0 else common_dir
    
    ppt_name += '.pptx'
    
    ppt_dir = os.path.join(configs.path1, ppt_name)
    ppt = pptx.Presentation()
    ppt.slide_height = 6858000    #设置ppt的高度
    ppt.slide_width = 12192000    #设置ppt的宽度
    pptw = pptx.util.Inches(configs.pptw) 
    ppth = pptx.util.Inches(configs.ppth) 
    pptt = pptx.util.Inches(configs.pptt) 
    pptg = pptx.util.Inches(configs.pptg) 
    #Mean
    
    #Setting Directory
    #Add the specific part of the directory name based on the filter
    common_dir = 'MeanCCMaps2_Nofilter' if configs.filter == 0.5 else 'MeanCCMaps2'

    #Add the specific part of the directory name based on the GSindex
    common_dir = common_dir + '_NGSR' if configs.GSindex == 0 else common_dir
    
    if configs.Taskindex == 1:
    #Add the specific part of the directory name based on the AtlasbasedROIregres
        common_dir = common_dir + '_TaskRegres' if configs.AtlasbasedROIregres == 1 else common_dir + '_TaskRegres36pixels'

    #Add the specific part of the directory name based on the Motionregres
    common_dir = common_dir + '_MotionReg' if configs.Motionregres == 1 else common_dir
    
    #Add the specific part of the directory name based on the bandwidth
    dir_name = common_dir + '_{}-{}Hz'.format(configs.fcs[configs.bandwidth][0], configs.fcs[configs.bandwidth][1]) if configs.bandwidth in [0,1,2,3,4] else common_dir
    
    #Add the specific part of the directory name based on the smoothing
    common_dir = common_dir + '_smooth' if configs.gaussian_sigma > 0 else common_dir
    
    MeanCCmapSlide = ppt.slides.add_slide(ppt.slide_layouts[6])
    add_text(MeanCCmapSlide, 'Mean of all trials' + dir_name, pptw, 0, pptw, ppth)
    for net_idx in range(len(configs.netname)):
        file_name = 'MeanCCmap_{}.tif'.format(configs.netname[net_idx]) if configs.gaussian_sigma == 0 else 'MeanCCmap_{}_smooth.tif'.format(configs.netname[net_idx])
        img_path = os.path.join(configs.path1, dir_name, file_name)
        
        textbox = add_text(MeanCCmapSlide, configs.netname[net_idx], (net_idx % 2) *(pptw + pptt), (np.floor(net_idx/2)+1)*(ppth+pptg), pptt, ppth)
        textbox.text_frame.paragraphs[0].font.size = pptx.util.Pt(10)
        
        MeanCCmapSlide.shapes.add_picture(img_path, (net_idx % 2) *(pptw+pptt)+ pptt, (np.floor(net_idx/2)+1)*(ppth+pptg), width=pptw, height=ppth)
            
    # Single scans
    
    #setting directory
    #Add the specific part of the directory name based on the GSindex
    common_dir = 'CCMaps' if configs.GSindex == 1 else 'CCMaps_NGSR'
    #Add the specific part of the directory name based on the filter
    if configs.filter == 0.5:
        common_dir = os.path.join('NoFilter', common_dir)
        
    #Add the specific part of the directory name based on the AtlasbasedROIregres
    if configs.Taskindex == 1:
        common_dir = os.path.join(common_dir, 'TaskRegres') if configs.AtlasbasedROIregres == 1 else os.path.join(common_dir, 'TaskRegres36pixels')
    #Add the specific part of the directory name based on the Motionregres
    if configs.Motionregres == 1:
        common_dir = os.path.join(common_dir, 'MotionReg')
    #Add the specific part of the directory name based on the bandwidth   
    common_dir = os.path.join(common_dir, '{}-{}Hz'.format(configs.fcs[configs.bandwidth][0], configs.fcs[configs.bandwidth][1])) if configs.bandwidth in [0,1,2,3,4] else common_dir
    #Add the specific part of the directory name based on the smoothing
    if configs.gaussian_sigma > 0:
        common_dir = os.path.join(common_dir, 'smooth')
        
    dir_name = common_dir
    
    for net_idx in range(len(configs.netname)):
        SingleScanSlide = ppt.slides.add_slide(ppt.slide_layouts[6])
        add_text(SingleScanSlide, configs.netname[net_idx], pptw+pptt, 0, pptw, ppth)
        

        for iifile_idx in range(len(configs.iifiles)):
            img_path = os.path.join(configs.path1, '{}'.format(configs.iifiles[iifile_idx]), dir_name, configs.netname[net_idx]+'.tif')
            
            textbox = add_text(SingleScanSlide, 'Scan \n#{}'.format(configs.iifiles[iifile_idx]), (iifile_idx % 2) * (pptw+pptt), (np.floor(iifile_idx/2)+1)*(ppth+pptg), pptt, ppth)
            textbox.text_frame.paragraphs[0].font.size = pptx.util.Pt(10)
            
            SingleScanSlide.shapes.add_picture(img_path, (iifile_idx % 2) * (pptw+pptt) + pptt, (np.floor(iifile_idx/2)+1)*(ppth+pptg), width=pptw, height=ppth)
                
    ppt.save(ppt_dir)